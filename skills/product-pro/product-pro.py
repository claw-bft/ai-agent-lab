#!/usr/bin/env python3
"""
Product Pro 核心实现
产品经理专业技能包 - 市场洞察、产品落地、数据驱动决策
"""

import os
import sys
import json
import argparse
import re
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Any
from dataclasses import dataclass, asdict

# 尝试导入依赖库
try:
    import requests
    REQUESTS_AVAILABLE = True
except ImportError:
    REQUESTS_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    # 定义占位类以避免类型注解错误
    class Presentation:
        pass

@dataclass
class CompetitorInfo:
    """竞品信息"""
    name: str
    positioning: str
    strengths: List[str]
    weaknesses: List[str]
    target_users: str
    pricing: str
    market_share: str = "未知"
    key_features: List[str] = None
    
    def __post_init__(self):
        if self.key_features is None:
            self.key_features = []

@dataclass
class PRDSection:
    """PRD文档章节"""
    title: str
    content: str
    subsections: List[Dict] = None
    
    def __post_init__(self):
        if self.subsections is None:
            self.subsections = []

class CompetitorAnalyzer:
    """竞品分析引擎"""
    
    def __init__(self):
        self.analysis_framework = {
            "swot": ["优势", "劣势", "机会", "威胁"],
            "4p": ["产品", "价格", "渠道", "推广"],
            "user_journey": ["认知", "考虑", "购买", "使用", "忠诚"]
        }
    
    def analyze(self, product_name: str, competitors: List[str] = None) -> Dict[str, Any]:
        """执行竞品分析"""
        
        if competitors is None:
            competitors = self._infer_competitors(product_name)
        
        analysis = {
            "success": True,
            "product": product_name,
            "timestamp": datetime.now().isoformat(),
            "market_overview": self._generate_market_overview(product_name),
            "competitors": [],
            "comparison_matrix": {},
            "insights": []
        }
        
        # 分析每个竞品
        for comp_name in competitors:
            comp_info = self._analyze_competitor(comp_name, product_name)
            analysis["competitors"].append(asdict(comp_info))
        
        # 生成对比矩阵
        analysis["comparison_matrix"] = self._generate_comparison_matrix(
            product_name, analysis["competitors"]
        )
        
        # 生成洞察
        analysis["insights"] = self._generate_insights(analysis["competitors"])
        
        return analysis
    
    def _infer_competitors(self, product_name: str) -> List[str]:
        """根据产品名称推断竞品"""
        product_lower = product_name.lower()
        
        competitor_map = {
            "ai代码": ["GitHub Copilot", "Cursor", "Codeium", "Tabnine"],
            "代码助手": ["GitHub Copilot", "Cursor", "Codeium", "Tabnine"],
            "笔记": ["Notion", "Obsidian", "Roam Research", "Logseq"],
            "协作": ["Slack", "Microsoft Teams", "飞书", "钉钉"],
            "云盘": ["Dropbox", "Google Drive", "OneDrive", "iCloud"],
            "设计": ["Figma", "Sketch", "Adobe XD", "Canva"],
            "项目管理": ["Jira", "Trello", "Asana", "Monday.com"],
            "文档": ["Google Docs", "Microsoft Word", "Notion", "Confluence"],
        }
        
        for key, comps in competitor_map.items():
            if key in product_lower:
                return comps
        
        # 默认返回通用竞品
        return ["市场领导者A", "创新者B", "价格竞争者C"]
    
    def _analyze_competitor(self, name: str, target_product: str) -> CompetitorInfo:
        """分析单个竞品"""
        # 这里可以接入真实搜索API获取数据
        # 目前使用模拟数据框架
        
        return CompetitorInfo(
            name=name,
            positioning=f"{name}在{target_product}领域的定位",
            strengths=[
                "功能完整度高",
                "用户基础大",
                "品牌知名度高"
            ],
            weaknesses=[
                "价格较高",
                "学习曲线陡峭",
                "定制化能力有限"
            ],
            target_users="企业用户和专业开发者",
            pricing="订阅制，$10-50/月",
            market_share="25-30%",
            key_features=["核心功能A", "核心功能B", "核心功能C"]
        )
    
    def _generate_market_overview(self, product_name: str) -> Dict:
        """生成市场概览"""
        return {
            "market_size": "预计市场规模数据（建议接入真实数据源）",
            "growth_rate": "年增长率约15-25%",
            "key_trends": [
                "AI集成成为标配",
                "协作功能日益重要",
                "移动端体验优化",
                "API生态建设"
            ],
            "barriers": [
                "技术门槛",
                "用户习惯",
                "数据安全合规"
            ]
        }
    
    def _generate_comparison_matrix(self, product_name: str, competitors: List[Dict]) -> Dict:
        """生成对比矩阵"""
        dimensions = ["功能完整度", "易用性", "价格", "性能", "支持服务"]
        
        matrix = {
            "dimensions": dimensions,
            "products": []
        }
        
        # 添加自身产品
        matrix["products"].append({
            "name": product_name,
            "scores": {d: "待评估" for d in dimensions}
        })
        
        # 添加竞品
        for comp in competitors:
            matrix["products"].append({
                "name": comp["name"],
                "scores": {d: "待评估" for d in dimensions}
            })
        
        return matrix
    
    def _generate_insights(self, competitors: List[Dict]) -> List[str]:
        """生成洞察建议"""
        return [
            "市场存在差异化机会，建议聚焦特定用户群体",
            "竞品在移动端体验方面普遍较弱，可作为突破口",
            "定价策略建议采用freemium模式降低用户门槛",
            "AI功能已成为行业标配，需尽快布局"
        ]

class PRDGenerator:
    """PRD文档生成器"""
    
    def __init__(self):
        self.templates = {
            "standard": self._standard_template,
            "lean": self._lean_template,
            "detailed": self._detailed_template
        }
    
    def generate(self, feature_name: str, template: str = "standard", 
                 context: Dict = None) -> Dict[str, Any]:
        """生成PRD文档"""
        
        if template not in self.templates:
            return {
                "success": False,
                "error": f"未知模板: {template}"
            }
        
        template_func = self.templates[template]
        sections = template_func(feature_name, context or {})
        
        return {
            "success": True,
            "feature": feature_name,
            "template": template,
            "timestamp": datetime.now().isoformat(),
            "sections": [asdict(s) for s in sections],
            "markdown": self._to_markdown(feature_name, sections)
        }
    
    def _standard_template(self, feature_name: str, context: Dict) -> List[PRDSection]:
        """标准PRD模板"""
        return [
            PRDSection(
                title="1. 文档概述",
                content=f"""### 1.1 文档信息
- **功能名称**: {feature_name}
- **版本**: 1.0
- **创建日期**: {datetime.now().strftime("%Y-%m-%d")}
- **状态**: 草稿

### 1.2 修订记录
| 版本 | 日期 | 修订人 | 修订内容 |
|------|------|--------|----------|
| 1.0 | {datetime.now().strftime("%Y-%m-%d")} | PM | 初稿创建 |

### 1.3 术语表
| 术语 | 定义 |
|------|------|
| {feature_name} | 本文档描述的核心功能 |
"""
            ),
            PRDSection(
                title="2. 背景与目标",
                content=f"""### 2.1 项目背景
[描述该功能产生的背景，解决的问题]

### 2.2 目标用户
- **主要用户**: [描述]
- **次要用户**: [描述]

### 2.3 业务目标
1. [目标1]
2. [目标2]
3. [目标3]

### 2.4 成功指标
| 指标 | 目标值 | 衡量方式 |
|------|--------|----------|
| 指标1 | 数值 | 方法 |
| 指标2 | 数值 | 方法 |
"""
            ),
            PRDSection(
                title="3. 功能需求",
                content=f"""### 3.1 功能清单

#### 3.1.1 {feature_name} - 核心功能
**需求描述**: 
[详细描述功能需求]

**验收标准**:
- [ ] 验收标准1
- [ ] 验收标准2
- [ ] 验收标准3

**优先级**: P0

#### 3.1.2 辅助功能
[其他相关功能]

### 3.2 用户故事
作为[用户角色]，我希望[需求]，以便[价值]

### 3.3 流程图
```
[用户流程图]
```
"""
            ),
            PRDSection(
                title="4. 非功能需求",
                content="""### 4.1 性能要求
- 页面加载时间 < 2秒
- API响应时间 < 500ms

### 4.2 兼容性
- 浏览器: Chrome, Firefox, Safari, Edge (最新2个版本)
- 移动端: iOS 14+, Android 10+

### 4.3 安全要求
- 用户数据加密传输
- 敏感操作需二次验证

### 4.4 可用性
- 系统可用性 > 99.9%
"""
            ),
            PRDSection(
                title="5. 界面原型",
                content="""### 5.1 页面结构
[页面布局描述]

### 5.2 交互说明
[关键交互说明]

### 5.3 视觉参考
[设计稿链接或描述]
"""
            ),
            PRDSection(
                title="6. 数据分析",
                content="""### 6.1 埋点需求
| 事件名 | 触发时机 | 属性 |
|--------|----------|------|
| event_name | 触发条件 | {属性} |

### 6.2 数据看板
[需要监控的数据指标]
"""
            ),
            PRDSection(
                title="7. 风险评估",
                content="""### 7.1 技术风险
| 风险 | 影响 | 缓解措施 |
|------|------|----------|
| 风险1 | 高/中/低 | 措施 |

### 7.2 业务风险
[业务层面风险]

### 7.3 合规风险
[合规相关考虑]
"""
            ),
            PRDSection(
                title="8. 发布计划",
                content="""### 8.1 里程碑
| 阶段 | 日期 | 交付物 |
|------|------|--------|
| 设计 | 日期 | 设计稿 |
| 开发 | 日期 | 功能代码 |
| 测试 | 日期 | 测试报告 |
| 上线 | 日期 | 正式发布 |

### 8.2 依赖项
- [ ] 依赖1
- [ ] 依赖2
"""
            )
        ]
    
    def _lean_template(self, feature_name: str, context: Dict) -> List[PRDSection]:
        """精简PRD模板"""
        return [
            PRDSection(
                title="1. 一句话描述",
                content=f"我们要做一个{feature_name}，帮助[目标用户]解决[问题]，实现[价值]。"
            ),
            PRDSection(
                title="2. 为什么做",
                content="""### 问题
[描述要解决的问题]

### 机会
[市场/用户机会]

### 不做会怎样
[不做的后果]
"""
            ),
            PRDSection(
                title="3. 怎么做",
                content="""### 核心功能
1. [功能点1]
2. [功能点2]

### 用户流程
[简化流程]

### 成功标准
- [标准1]
- [标准2]
"""
            ),
            PRDSection(
                title="4. 关键指标",
                content="""| 指标 | 当前 | 目标 |
|------|------|------|
| 指标1 | 数值 | 数值 |
"""
            )
        ]
    
    def _detailed_template(self, feature_name: str, context: Dict) -> List[PRDSection]:
        """详细PRD模板"""
        # 在标准模板基础上增加更多细节
        sections = self._standard_template(feature_name, context)
        
        # 增加技术实现章节
        sections.insert(4, PRDSection(
            title="4. 技术方案",
            content="""### 4.1 架构设计
[系统架构图]

### 4.2 接口设计
#### API列表
| 接口 | 方法 | 描述 |
|------|------|------|
| /api/xxx | GET | 描述 |

#### 请求/响应示例
```json
{
  "code": 200,
  "data": {}
}
```

### 4.3 数据库设计
[ER图或表结构]

### 4.4 第三方依赖
- [依赖1]: 用途
- [依赖2]: 用途
"""
        ))
        
        return sections
    
    def _to_markdown(self, feature_name: str, sections: List[PRDSection]) -> str:
        """转换为Markdown格式"""
        md = f"# {feature_name} - 产品需求文档\n\n"
        md += f"> 生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n\n"
        md += "---\n\n"
        
        for section in sections:
            md += f"\n{section.content}\n\n"
            md += "---\n"
        
        return md

class PPTGenerator:
    """PPT生成器"""
    
    def __init__(self):
        self.slide_layouts = {
            "title": self._create_title_slide,
            "content": self._create_content_slide,
            "bullet": self._create_bullet_slide,
            "two_column": self._create_two_column_slide,
            "image": self._create_image_slide
        }
    
    def generate(self, topic: str, slides_count: int = 10, 
                 outline: List[Dict] = None) -> Dict[str, Any]:
        """生成PPT"""
        
        if not PPTX_AVAILABLE:
            return {
                "success": False,
                "error": "python-pptx未安装，请运行: pip install python-pptx"
            }
        
        if outline is None:
            outline = self._generate_default_outline(topic, slides_count)
        
        try:
            prs = Presentation()
            
            for slide_info in outline:
                layout_type = slide_info.get("type", "content")
                if layout_type in self.slide_layouts:
                    self.slide_layouts[layout_type](prs, slide_info)
            
            # 保存文件
            output_path = f"{topic.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d')}.pptx"
            prs.save(output_path)
            
            return {
                "success": True,
                "topic": topic,
                "slides_count": len(outline),
                "output_path": output_path,
                "outline": outline
            }
        
        except Exception as e:
            return {
                "success": False,
                "error": f"生成PPT失败: {str(e)}"
            }
    
    def _generate_default_outline(self, topic: str, count: int) -> List[Dict]:
        """生成默认大纲"""
        outline = [
            {
                "type": "title",
                "title": topic,
                "subtitle": f"产品规划报告 - {datetime.now().strftime('%Y年%m月')}"
            },
            {
                "type": "bullet",
                "title": "目录",
                "items": [
                    "市场分析",
                    "产品定位",
                    "核心功能",
                    "发展规划",
                    "预期收益"
                ]
            },
            {
                "type": "content",
                "title": "市场分析",
                "content": "[市场规模、增长趋势、竞争格局分析]"
            },
            {
                "type": "two_column",
                "title": "目标用户",
                "left": {
                    "title": "核心用户",
                    "items": ["特征1", "特征2", "特征3"]
                },
                "right": {
                    "title": "用户需求",
                    "items": ["需求1", "需求2", "需求3"]
                }
            },
            {
                "type": "bullet",
                "title": "产品定位",
                "items": [
                    "核心价值主张",
                    "差异化优势",
                    "竞争策略"
                ]
            },
            {
                "type": "content",
                "title": "核心功能",
                "content": "[主要功能模块介绍]"
            },
            {
                "type": "two_column",
                "title": "发展规划",
                "left": {
                    "title": "短期目标 (3个月)",
                    "items": ["MVP发布", "种子用户获取", "核心功能验证"]
                },
                "right": {
                    "title": "长期愿景 (12个月)",
                    "items": ["规模化增长", "生态建设", "市场领导地位"]
                }
            },
            {
                "type": "bullet",
                "title": "关键指标",
                "items": [
                    "用户增长: 月活达到X万",
                    "留存率: 次日留存X%",
                    "收入: 月收入X万元"
                ]
            },
            {
                "type": "content",
                "title": "资源需求",
                "content": "[人力、资金、时间等资源需求]"
            },
            {
                "type": "title",
                "title": "谢谢",
                "subtitle": "Questions?"
            }
        ]
        
        return outline[:count]
    
    def _create_title_slide(self, prs: Presentation, info: Dict):
        """创建标题页"""
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = info.get("title", "")
        subtitle.text = info.get("subtitle", "")
    
    def _create_content_slide(self, prs: Presentation, info: Dict):
        """创建内容页"""
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = info.get("title", "")
        content.text = info.get("content", "")
    
    def _create_bullet_slide(self, prs: Presentation, info: Dict):
        """创建列表页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        body = slide.placeholders[1]
        
        title.text = info.get("title", "")
        
        tf = body.text_frame
        tf.text = info.get("items", [""])[0] if info.get("items") else ""
        
        for item in info.get("items", [])[1:]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
    
    def _create_two_column_slide(self, prs: Presentation, info: Dict):
        """创建双栏页"""
        slide_layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title_shape = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        title_shape.text_frame.text = info.get("title", "")
        
        # 左栏
        left = info.get("left", {})
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(4.5), Inches(5)
        )
        tf = left_box.text_frame
        tf.text = left.get("title", "")
        
        for item in left.get("items", []):
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.level = 1
        
        # 右栏
        right = info.get("right", {})
        right_box = slide.shapes.add_textbox(
            Inches(5), Inches(1.5), Inches(4.5), Inches(5)
        )
        tf = right_box.text_frame
        tf.text = right.get("title", "")
        
        for item in right.get("items", []):
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.level = 1
    
    def _create_image_slide(self, prs: Presentation, info: Dict):
        """创建图片页"""
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(1)
        )
        title.text_frame.text = info.get("title", "")
        
        # 图片占位
        content = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(9), Inches(5)
        )
        content.text_frame.text = "[图片占位符]"

class MarketResearcher:
    """市场研究模块"""
    
    def conduct_research(self, topic: str, method: str = "secondary",
                        target: str = None, sample_size: int = 0) -> Dict[str, Any]:
        """执行市场研究"""
        
        methods = {
            "secondary": self._secondary_research,
            "interview": self._interview_research,
            "survey": self._survey_research,
            "competitive": self._competitive_research
        }
        
        if method not in methods:
            return {
                "success": False,
                "error": f"未知研究方法: {method}"
            }
        
        return methods[method](topic, target, sample_size)
    
    def _secondary_research(self, topic: str, target: str, sample_size: int) -> Dict:
        """二手资料研究"""
        return {
            "success": True,
            "method": "secondary_research",
            "topic": topic,
            "findings": {
                "market_size": "[需接入真实数据源]",
                "trends": [
                    "趋势1: AI技术普及",
                    "趋势2: 用户体验优先",
                    "趋势3: 数据驱动决策"
                ],
                "key_players": ["公司A", "公司B", "公司C"],
                "data_sources": [
                    "行业报告",
                    "新闻资讯",
                    "学术论文"
                ]
            },
            "recommendations": [
                "建议接入Tavily等搜索API获取实时数据",
                "建立行业专家访谈名单"
            ]
        }
    
    def _interview_research(self, topic: str, target: str, sample_size: int) -> Dict:
        """用户访谈研究"""
        return {
            "success": True,
            "method": "user_interview",
            "topic": topic,
            "target_users": target or "未指定",
            "sample_size": sample_size or 10,
            "interview_guide": {
                "opening": ["请介绍一下您的背景"],
                "main": [
                    f"您在使用{topic}时遇到的最大挑战是什么？",
                    "您目前的解决方案是什么？",
                    "您理想中的解决方案是怎样的？"
                ],
                "closing": ["还有什么想补充的吗？"]
            },
            "note_template": """
## 访谈记录模板

### 基本信息
- 受访者:
- 时间:
- 时长:

### 关键洞察
1.
2.
3.

### 痛点总结
-

### 机会点
-
"""
        }
    
    def _survey_research(self, topic: str, target: str, sample_size: int) -> Dict:
        """问卷调查"""
        return {
            "success": True,
            "method": "survey",
            "topic": topic,
            "questionnaire": {
                "screening": [
                    "您是否使用过相关产品？",
                    "您的职业是？"
                ],
                "main_questions": [
                    {
                        "question": f"您对{topic}的满意度如何？",
                        "type": "scale",
                        "options": ["1-非常不满意", "5-非常满意"]
                    },
                    {
                        "question": "您最看重的功能是什么？",
                        "type": "multi_select",
                        "options": ["功能A", "功能B", "功能C", "其他"]
                    }
                ],
                "demographics": [
                    "年龄段",
                    "所在城市",
                    "收入水平"
                ]
            },
            "target_responses": sample_size or 100
        }
    
    def _competitive_research(self, topic: str, target: str, sample_size: int) -> Dict:
        """竞品研究"""
        analyzer = CompetitorAnalyzer()
        return analyzer.analyze(topic)

def main():
    parser = argparse.ArgumentParser(description="Product Pro - 产品经理专业技能包")
    subparsers = parser.add_subparsers(dest="command", help="可用命令")
    
    # 竞品分析命令
    competitor_parser = subparsers.add_parser("competitor", help="竞品分析")
    competitor_parser.add_argument("action", choices=["analyze"])
    competitor_parser.add_argument("--product", required=True, help="产品名称")
    competitor_parser.add_argument("--competitors", help="竞品列表 (逗号分隔)")
    competitor_parser.add_argument("--output", help="输出文件路径")
    competitor_parser.add_argument("--json", action="store_true", help="JSON格式输出")
    
    # PRD生成命令
    prd_parser = subparsers.add_parser("prd", help="PRD文档生成")
    prd_parser.add_argument("action", choices=["create"])
    prd_parser.add_argument("--feature", required=True, help="功能名称")
    prd_parser.add_argument("--template", default="standard", 
                           choices=["standard", "lean", "detailed"],
                           help="PRD模板类型")
    prd_parser.add_argument("--output", help="输出文件路径")
    prd_parser.add_argument("--json", action="store_true", help="JSON格式输出")
    
    # PPT生成命令
    ppt_parser = subparsers.add_parser("ppt", help="PPT生成")
    ppt_parser.add_argument("action", choices=["create"])
    ppt_parser.add_argument("--topic", required=True, help="PPT主题")
    ppt_parser.add_argument("--slides", type=int, default=10, help="幻灯片数量")
    ppt_parser.add_argument("--output", help="输出文件路径")
    ppt_parser.add_argument("--json", action="store_true", help="JSON格式输出")
    
    # 市场研究命令
    research_parser = subparsers.add_parser("research", help="市场研究")
    research_parser.add_argument("action", choices=["conduct"])
    research_parser.add_argument("--topic", required=True, help="研究主题")
    research_parser.add_argument("--method", default="secondary",
                                choices=["secondary", "interview", "survey", "competitive"],
                                help="研究方法")
    research_parser.add_argument("--target", help="目标用户群体")
    research_parser.add_argument("--n", type=int, default=0, help="样本数量")
    research_parser.add_argument("--json", action="store_true", help="JSON格式输出")
    
    args = parser.parse_args()
    
    if not args.command:
        parser.print_help()
        sys.exit(1)
    
    result = None
    
    if args.command == "competitor":
        analyzer = CompetitorAnalyzer()
        competitors = None
        if args.competitors:
            competitors = [c.strip() for c in args.competitors.split(",")]
        result = analyzer.analyze(args.product, competitors)
    
    elif args.command == "prd":
        generator = PRDGenerator()
        result = generator.generate(args.feature, args.template)
        
        # 保存Markdown文件
        if result.get("success") and not args.json:
            output_path = args.output or f"PRD_{args.feature.replace(' ', '_')}.md"
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(result.get("markdown", ""))
            result["saved_to"] = output_path
    
    elif args.command == "ppt":
        generator = PPTGenerator()
        result = generator.generate(args.topic, args.slides)
    
    elif args.command == "research":
        researcher = MarketResearcher()
        result = researcher.conduct_research(
            args.topic, args.method, args.target, args.n
        )
    
    # 输出结果
    if args.json or (result and not result.get("success")):
        print(json.dumps(result, indent=2, ensure_ascii=False))
    else:
        if result.get("success"):
            print(f"✓ {args.command} 执行成功")
            if "saved_to" in result:
                print(f"  文件已保存: {result['saved_to']}")
            if "output_path" in result:
                print(f"  文件已生成: {result['output_path']}")
        else:
            print(f"✗ 错误: {result.get('error', '未知错误')}")
            sys.exit(1)

if __name__ == "__main__":
    main()
